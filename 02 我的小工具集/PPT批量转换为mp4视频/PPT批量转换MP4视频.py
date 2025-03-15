from pptx import Presentation
from moviepy import ImageSequenceClip
import os

def ppt_to_video(ppt_path, video_path):
    # 读取PPT文件
    prs = Presentation(ppt_path)
    images = []
    
    # 遍历每一页幻灯片
    for slide in prs.slides:
        slide_path = f"slide_{prs.slides.index(slide)}.png"
        slide.save(slide_path)
        images.append(slide_path)
    
    # 使用MoviePy生成视频
    clip = ImageSequenceClip(images, fps=1)
    clip.write_videofile(video_path, codec='libx264')
    
    # 清理临时图片
    for img in images:
        os.remove(img)

# 只转换一个文件
ppt_folder = "/Users/harmnet/Desktop/PPT"
video_folder = "/Users/harmnet/Desktop/PPT视频"
os.makedirs(video_folder, exist_ok=True)

# 获取文件夹中的第一个PPTX文件
ppt_files = [f for f in os.listdir(ppt_folder) if f.endswith(".pptx")]
if ppt_files:
    ppt_file = ppt_files[0]  # 只取第一个文件
    print(f"正在转换: {ppt_file}")
    ppt_path = os.path.join(ppt_folder, ppt_file)
    video_path = os.path.join(video_folder, ppt_file.replace(".pptx", ".mp4"))
    ppt_to_video(ppt_path, video_path)
    print(f"转换完成: {video_path}")
else:
    print(f"在 {ppt_folder} 中没有找到PPTX文件")