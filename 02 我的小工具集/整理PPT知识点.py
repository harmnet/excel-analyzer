import os
import pandas as pd
from pptx import Presentation
import re
import requests
import json
import time
import sys
from tqdm.auto import tqdm  # 使用tqdm.auto以自动选择最佳的进度条类型

# 配置路径
desktop_path = "/Users/harmnet/Desktop"
ppt_folder_path = os.path.join(desktop_path, "02 PPT课件")
output_file = os.path.join(desktop_path, "知识结构树.xlsx")

# 硅基流动API配置
API_URL = "https://api.siliconflow.cn/v1/chat/completions"
API_KEY = "sk-bivnwauskdbvpspvmdorrgkrpwlyfxbfcezqsfsevowzubdj"  # 请替换为你的API密钥

# 初始化结果列表
results = []
serial_number = 1

# 提取PPT内容的函数
def extract_ppt_content(ppt_path):
    try:
        prs = Presentation(ppt_path)
        slides_content = []  # 保存每页PPT的内容
        
        for slide in prs.slides:
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            slides_content.append(slide_text)
        
        return slides_content  # 返回列表，每个元素是一页PPT的内容
    except Exception as e:
        print(f"处理PPT文件时出错: {ppt_path}, 错误: {str(e)}")
        return []

# 使用硅基流动API提取知识点
def extract_knowledge_points(content, chapter, section):
    try:
        prompt = f"""
        请分析以下PPT内容，提取出主要任务和每个任务对应的关键知识点。
        PPT章节: {chapter}
        PPT小节: {section}
        
        PPT内容:
        {content[:3000]}  # 限制内容长度以避免超出token限制
        
        请按以下格式输出:
        任务1: [任务名称]
        知识点: [知识点1]; [知识点2]; [知识点3]...
        
        任务2: [任务名称]
        知识点: [知识点1]; [知识点2]; [知识点3]...
        
        ...以此类推
        """
        
        headers = {
            "Authorization": f"Bearer {API_KEY}",
            "Content-Type": "application/json"
        }
        
        payload = {
            "model": "Pro/deepseek-ai/DeepSeek-V3",  # 正确的模型名称
            "messages": [
                {"role": "system", "content": "你是一个专业的PPT内容分析助手，擅长提取PPT中的任务和知识点。"},
                {"role": "user", "content": prompt}
            ],
            "stream": False,
            "max_tokens": 1000,
            "stop": ["null"],
            "temperature": 0.3,
            "top_p": 0.7,
            "top_k": 50,
            "frequency_penalty": 0.5,
            "n": 1,
            "response_format": {"type": "text"}
        }
        
        response = requests.post(API_URL, headers=headers, json=payload)  # 使用json参数替代data
        
        if response.status_code == 200:
            result = response.json()
            return result["choices"][0]["message"]["content"]
        else:
            print(f"API调用失败，状态码: {response.status_code}, 响应: {response.text}")
            return "无法提取知识点"
            
    except Exception as e:
        print(f"调用AI接口时出错: {str(e)}")
        time.sleep(5)  # 出错时等待一段时间再重试
        return "无法提取知识点"

# 解析AI返回的结果
def parse_ai_response(response):
    tasks = []
    
    # 使用正则表达式匹配任务和知识点
    pattern = r"任务\d+:\s*(.*?)\n知识点:\s*(.*?)(?=\n任务\d+:|$)"
    matches = re.finditer(pattern, response, re.DOTALL)
    
    for match in matches:
        task = match.group(1).strip()
        knowledge_points = match.group(2).strip()
        tasks.append((task, knowledge_points))
    
    # 如果没有匹配到任何任务，尝试其他可能的格式
    if not tasks:
        lines = response.split('\n')
        current_task = None
        current_points = []
        
        for line in lines:
            if line.strip().startswith("任务") or ":" in line and not line.lower().startswith("知识点"):
                if current_task:
                    tasks.append((current_task, "; ".join(current_points)))
                current_task = line.split(":", 1)[1].strip() if ":" in line else line.strip()
                current_points = []
            elif line.strip().startswith("知识点") or "知识点" in line.lower():
                points = line.split(":", 1)[1].strip() if ":" in line else line.strip()
                current_points.append(points)
            elif current_task and line.strip():
                current_points.append(line.strip())
        
        if current_task:
            tasks.append((current_task, "; ".join(current_points)))
    
    return tasks

# 批量处理PPT内容并调用API
def batch_process_ppt(slides_content, chapter, section, batch_size=10):
    all_tasks = []
    total_slides = len(slides_content)
    batch_count = (total_slides + batch_size - 1) // batch_size  # 向上取整
    
    print(f"PPT共有{total_slides}页，分为{batch_count}批处理")
    
    # 使用tqdm显示批次处理进度
    batch_pbar = tqdm(range(batch_count), desc="批次进度", leave=False)
    for i in batch_pbar:
        start_idx = i * batch_size
        end_idx = min((i + 1) * batch_size, total_slides)
        batch_slides = slides_content[start_idx:end_idx]
        
        # 更新进度条描述
        batch_pbar.set_description(f"处理页码 {start_idx+1}-{end_idx}")
        
        # 合并当前批次的内容
        batch_content = "\n\n--- 第" + "、".join([str(j+1) for j in range(start_idx, end_idx)]) + "页 ---\n\n"
        batch_content += "\n\n".join(batch_slides)
        
        # 定义批处理提示词
        batch_prompt = f"""
        请分析以下PPT内容（第{start_idx+1}页到第{end_idx}页），提取出主要任务和每个任务对应的关键知识点。知识点请尽量简洁的总结出来。
        PPT章节: {chapter}
        PPT小节: {section}

        PPT内容:
        {batch_content}

        请按以下格式输出:
        任务1: [任务名称]
        知识点: [知识点1]; [知识点2]; [知识点3]...

        任务2: [任务名称]
        知识点: [知识点1]; [知识点2]; [知识点3]...

        ...以此类推
        """
        
        # 调用AI提取知识点
        batch_pbar.set_postfix(status="调用API中...")
        ai_response = extract_knowledge_points(batch_prompt, chapter, section)
        batch_pbar.set_postfix(status="解析结果中...")
        batch_tasks = parse_ai_response(ai_response)
        
        if batch_tasks:
            # 给每个任务添加批次信息
            batch_tasks_with_info = []
            for task, knowledge_points in batch_tasks:
                task_with_page = f"{task} (页码: {start_idx+1}-{end_idx})"
                batch_tasks_with_info.append((task_with_page, knowledge_points))
            all_tasks.extend(batch_tasks_with_info)
            batch_pbar.set_postfix(status=f"提取到{len(batch_tasks)}个任务")
        else:
            batch_pbar.set_postfix(status="未提取到任务")
        
        # 避免API调用过于频繁
        if i < batch_count - 1:  # 不是最后一批
            batch_pbar.set_postfix(status="等待中...")
            time.sleep(5)
    
    # 如果没有提取到任何任务，返回默认值
    if not all_tasks:
        all_tasks = [("未识别到任务", "未识别到知识点")]
    
    return all_tasks

# 主函数
def main():
    global serial_number
    
    if not os.path.exists(ppt_folder_path):
        print(f"文件夹不存在: {ppt_folder_path}")
        return
    
    # 获取所有章节文件夹
    chapter_folders = [f for f in os.listdir(ppt_folder_path) if os.path.isdir(os.path.join(ppt_folder_path, f)) and not f.startswith('.')]
    total_chapters = len(chapter_folders)
    print(f"总共发现 {total_chapters} 个章节文件夹")
    
    # 计算总文件数以显示整体进度
    total_files = 0
    for chapter_folder in chapter_folders:
        chapter_path = os.path.join(ppt_folder_path, chapter_folder)
        ppt_files = [f for f in os.listdir(chapter_path) if f.endswith('.pptx') and not f.startswith('~$')]
        total_files += len(ppt_files)
    
    print(f"总共需要处理 {total_files} 个PPT文件")
    
    # 使用两个嵌套的进度条：外层为总体进度，内层为当前章节进度
    processed_files = 0
    overall_pbar = tqdm(total=total_files, desc="总体进度", position=0)
    
    # 遍历文件夹
    for chapter_idx, chapter_folder in enumerate(chapter_folders, 1):
        chapter_path = os.path.join(ppt_folder_path, chapter_folder)
        
        # 遍历章节文件夹中的PPT文件
        ppt_files = [f for f in os.listdir(chapter_path) if f.endswith('.pptx') and not f.startswith('~$')]
        
        # 更新章节进度信息
        chapter_desc = f"章节 [{chapter_idx}/{total_chapters}]: {chapter_folder}"
        print(f"\n{chapter_desc}")
        
        # 章节进度条
        chapter_pbar = tqdm(total=len(ppt_files), desc=chapter_desc, position=1, leave=False)
        
        for ppt_idx, ppt_file in enumerate(ppt_files, 1):
            ppt_path = os.path.join(chapter_path, ppt_file)
            section_name = os.path.splitext(ppt_file)[0]  # 去掉文件扩展名
            
            processed_files += 1
            file_desc = f"文件 [{processed_files}/{total_files}]: {ppt_file}"
            
            # 更新进度条描述
            chapter_pbar.set_description(file_desc)
            overall_pbar.set_description(f"总进度: {processed_files}/{total_files} ({(processed_files/total_files*100):.1f}%)")
            
            # 提取PPT内容（获取每页的内容列表）
            chapter_pbar.set_postfix(status="提取PPT内容...")
            slides_content = extract_ppt_content(ppt_path)
            
            if not slides_content:
                chapter_pbar.set_postfix(status="内容为空，跳过")
                chapter_pbar.update(1)
                overall_pbar.update(1)
                continue
            
            # 批量处理PPT内容并调用AI提取知识点
            chapter_pbar.set_postfix(status="提取知识点...")
            tasks = batch_process_ppt(slides_content, chapter_folder, section_name)
            
            # 如果没有提取到任务，添加一个默认任务
            if not tasks or tasks[0][0] == "未识别到任务":
                results.append({
                    "序号": serial_number,
                    "章": chapter_folder,
                    "小节": section_name,
                    "任务": "未识别到任务",
                    "知识点": "未识别到知识点"
                })
                serial_number += 1
                chapter_pbar.set_postfix(status="未提取到任务")
            else:
                # 添加每个任务和知识点到结果列表
                for task, knowledge_points in tasks:
                    results.append({
                        "序号": serial_number,
                        "章": chapter_folder,
                        "小节": section_name,
                        "任务": task, 
                        "知识点": knowledge_points
                    })
                    serial_number += 1
                chapter_pbar.set_postfix(status=f"提取到{len(tasks)}个任务")
            
            # 更新进度条
            chapter_pbar.update(1)
            overall_pbar.update(1)
            
            # 每处理完一个文件暂停一下，避免API调用过于频繁
            time.sleep(1)
        
        # 关闭章节进度条
        chapter_pbar.close()
        print(f"章节 {chapter_folder} 处理完成 ({chapter_idx}/{total_chapters})")
    
    # 关闭总体进度条
    overall_pbar.close()
    print(f"所有文件处理完成！总共处理了 {processed_files} 个文件")
    
    # 将结果保存为Excel文件
    if results:
        print("正在保存结果到Excel...")
        
        # 处理知识点，将每个知识点单独一行
        expanded_results = []
        for result in results:
            # 检查是否有知识点
            if result["知识点"] != "未识别到知识点":
                # 使用分号分割知识点
                knowledge_points = result["知识点"].split(";")
                # 清理每个知识点
                knowledge_points = [kp.strip() for kp in knowledge_points if kp.strip()]
                
                # 为每个知识点创建一行
                for i, kp in enumerate(knowledge_points):
                    new_row = result.copy()
                    # 修改知识点为单个知识点
                    new_row["知识点"] = kp
                    
                    # 如果不是第一个知识点，修改任务名称添加标记
                    if i > 0:
                        new_row["任务"] = f"{result['任务']} (知识点 {i+1})"
                    
                    expanded_results.append(new_row)
            else:
                # 如果没有识别到知识点，保持原样
                expanded_results.append(result)
        
        # 创建DataFrame并保存
        df = pd.DataFrame(expanded_results)
        df.to_excel(output_file, index=False)
        print(f"知识结构树已保存到: {output_file}")
    else:
        print("未提取到任何知识点")

if __name__ == "__main__":
    main()
